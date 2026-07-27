import os
import time
from io import BytesIO
from pathlib import Path
from datetime import date
from datetime import datetime, timezone

from fastapi import Depends, FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from sqlalchemy.orm import Session

from .database import Base, database_backend, engine, get_db
from .models import AgentFeedback, AgentInteraction, CampusRequest, InventoryItem, KnowledgeDocument, Sale
from .schemas import (
    ChatRequest,
    ChatResponse,
    DataCreateResponse,
    DocumentCreate,
    ForecastPointCreate,
    FeedbackCreate,
    InventoryCreate,
    RequestAction,
    SaleCreate,
    UploadResponse,
)
from .services.analytics import dashboard_summary
from .services.forecast import forecast_next_month
from .services.gemini import generate_grounded_answer
from .services.intent_classifier import classifier_evaluation, classify_inquiry
from .services.seed import ensure_specialized_knowledge, seed_database
from .services.vector_store import index_document, index_missing_documents, initialize_vector_store, vector_store_status
from .services.workflow import run_agent_workflow


DEFAULT_ALLOWED_ORIGINS = [
    "http://localhost:5173",
    "http://127.0.0.1:5173",
    "https://school-ai-campus-assistant.vercel.app",
]


def allowed_origins() -> list[str]:
    configured = os.getenv("FRONTEND_ORIGINS", "")
    extra_origins = [origin.strip().rstrip("/") for origin in configured.split(",") if origin.strip()]
    return [*DEFAULT_ALLOWED_ORIGINS, *extra_origins]


app = FastAPI(
    title="Ari Campus AI",
    description="Agentic AI assistant for school records, document search, reports, and activity forecasting.",
    version="1.0.0",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins(),
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.on_event("startup")
def startup() -> None:
    Base.metadata.create_all(bind=engine)
    db = next(get_db())
    try:
        seed_database(db)
        ensure_specialized_knowledge(db)
        initialize_vector_store(db)
    finally:
        db.close()


@app.get("/")
def api_home() -> dict:
    return {
        "name": "Ari Campus AI API",
        "status": "ok",
        "frontend": "https://school-ai-campus-assistant.vercel.app",
        "health": "/health",
        "docs": "/docs",
        "chat": "/chat",
    }


@app.get("/health")
def health() -> dict:
    return {
        "status": "ok",
        "gemini_configured": bool(os.getenv("GEMINI_API_KEY", "").strip()),
        "gemini_model": os.getenv("GEMINI_MODEL", "gemini-2.5-flash"),
        "database": database_backend(),
        "orchestration": "LangGraph",
        "retrieval": vector_store_status(),
    }


@app.get("/dashboard/summary")
def get_dashboard_summary(db: Session = Depends(get_db)) -> dict:
    return dashboard_summary(db)


@app.get("/forecast/sales")
def get_sales_forecast(db: Session = Depends(get_db)) -> dict:
    return forecast_next_month(db)


@app.post("/chat", response_model=ChatResponse)
def chat(request: ChatRequest, db: Session = Depends(get_db)) -> dict:
    started = time.perf_counter()
    result = run_agent_workflow(db, request.message, role=request.role, student_name=request.student_name)
    intent = result.get("intent", "")
    domain_agents = {
        "enrollment_request_agent": ("Ari · Enrollment Assistant", "Student request → Registrar approval queue"),
        "borrowing_request_agent": ("Ari · Equipment Assistant", "Inventory check → Admin approval queue"),
        "registrar_domain_agent": ("Registrar Agent", "Enrollment and registrar knowledge"),
        "finance_domain_agent": ("Finance Agent", "Scholarship and finance knowledge"),
        "guidance_domain_agent": ("Guidance Agent", "Student support knowledge"),
        "library_domain_agent": ("Library Agent", "Library services knowledge"),
        "it_domain_agent": ("IT Agent", "Campus technology knowledge"),
    }
    if intent in domain_agents:
        agent_name, data_path = domain_agents[intent]
    elif "inventory" in intent or "sql" in intent:
        agent_name, data_path = "Records Agent", "SQLAlchemy structured records"
    elif "forecast" in intent:
        agent_name, data_path = "Forecast Agent", "scikit-learn forecast service"
    elif "rag" in intent:
        agent_name, data_path = "Knowledge Agent", "RAG document retrieval"
    elif "report" in intent:
        agent_name, data_path = "Report Agent", "Records + RAG + forecast tools"
    else:
        agent_name, data_path = "Campus Router", "Intent router + campus tools"
    if result.get("sources"):
        retrieval_method = result["sources"][0].get("retrieval_method")
        if retrieval_method:
            data_path = retrieval_method
    result["agent_name"] = agent_name
    result["data_path"] = data_path
    generated = generate_grounded_answer(request.message, request.role, result)
    result["answer"] = generated["answer"]
    result["llm_provider"] = generated["provider"]
    result["llm_model"] = generated["model"]
    result["llm_status"] = generated["status"]
    result["response_time_ms"] = round((time.perf_counter() - started) * 1000)
    result.setdefault("model_trace", classify_inquiry(request.message))
    db.add(
        AgentInteraction(
            question=request.message,
            role=request.role,
            intent=intent,
            agent_name=agent_name,
            latency_ms=result["response_time_ms"],
            sources_used=len(result.get("sources", [])),
            successful=True,
            created_at=datetime.now(timezone.utc),
        )
    )
    db.commit()
    return result


@app.get("/ml/evaluation")
def get_ml_evaluation() -> dict:
    return classifier_evaluation()


@app.get("/admin/metrics")
def get_admin_metrics(db: Session = Depends(get_db)) -> dict:
    interactions = db.query(AgentInteraction).order_by(AgentInteraction.id.desc()).limit(500).all()
    feedback_count = db.query(AgentFeedback).count()
    open_tickets = db.query(AgentFeedback).filter(AgentFeedback.escalated.is_(True)).count()
    documents = db.query(KnowledgeDocument).count()
    total = len(interactions)
    average_latency = round(sum(item.latency_ms for item in interactions) / total) if total else 0
    source_grounded = sum(1 for item in interactions if item.sources_used > 0)
    agent_usage: dict[str, int] = {}
    for item in interactions:
        agent_usage[item.agent_name] = agent_usage.get(item.agent_name, 0) + 1
    return {
        "total_interactions": total,
        "average_latency_ms": average_latency,
        "source_grounded_answers": source_grounded,
        "feedback_count": feedback_count,
        "open_tickets": open_tickets,
        "indexed_documents": documents,
        "agent_usage": [{"agent": name, "requests": count} for name, count in sorted(agent_usage.items(), key=lambda pair: pair[1], reverse=True)],
        "recent": [
            {
                "question": item.question[:100],
                "role": item.role,
                "agent": item.agent_name,
                "latency_ms": item.latency_ms,
                "sources_used": item.sources_used,
            }
            for item in interactions[:8]
        ],
    }


def request_payload(item: CampusRequest) -> dict:
    return {
        "id": item.id,
        "request_type": item.request_type,
        "student_name": item.student_name,
        "sku": item.sku,
        "product": item.product,
        "quantity": item.quantity,
        "details": item.details,
        "status": item.status,
        "reviewed_by": item.reviewed_by,
        "created_at": item.created_at.isoformat(),
        "reviewed_at": item.reviewed_at.isoformat() if item.reviewed_at else None,
        "returned_at": item.returned_at.isoformat() if item.returned_at else None,
    }


@app.get("/requests")
def list_requests(role: str = "Admin", student_name: str | None = None, db: Session = Depends(get_db)) -> list[dict]:
    query = db.query(CampusRequest)
    if role.lower() == "student" and student_name:
        query = query.filter(CampusRequest.student_name == student_name)
    elif role.lower() == "registrar":
        query = query.filter(CampusRequest.request_type == "enrollment")
    elif role.lower() == "admin":
        query = query.filter(CampusRequest.request_type == "borrow")
    return [request_payload(item) for item in query.order_by(CampusRequest.id.desc()).limit(50).all()]


@app.post("/requests/{request_id}/approve", response_model=DataCreateResponse)
def approve_request(request_id: int, request: RequestAction, db: Session = Depends(get_db)) -> dict:
    item = db.query(CampusRequest).filter(CampusRequest.id == request_id).first()
    if not item:
        raise HTTPException(status_code=404, detail="Request not found.")
    if item.status not in {"pending_registrar", "pending_admin"}:
        raise HTTPException(status_code=409, detail="Request is no longer pending.")

    if item.request_type == "borrow":
        inventory = db.query(InventoryItem).filter(InventoryItem.sku == item.sku).first()
        if not inventory:
            raise HTTPException(status_code=409, detail="This item is no longer stocked by the school.")
        if inventory.stock_remaining < item.quantity:
            raise HTTPException(status_code=409, detail=f"Only {inventory.stock_remaining} item(s) remain in stock.")
        inventory.stock_remaining -= item.quantity
        item.status = "borrowed"
    else:
        item.status = "enrolled"

    item.reviewed_by = request.reviewer
    item.reviewed_at = datetime.now(timezone.utc)
    db.commit()
    db.refresh(item)
    return {"kind": item.request_type, "message": f"Request #{item.id} is now {item.status}.", "item": request_payload(item)}


@app.post("/requests/{request_id}/return", response_model=DataCreateResponse)
def return_borrowed_item(request_id: int, request: RequestAction, db: Session = Depends(get_db)) -> dict:
    item = db.query(CampusRequest).filter(CampusRequest.id == request_id).first()
    if not item or item.request_type != "borrow":
        raise HTTPException(status_code=404, detail="Borrowing request not found.")
    if item.status != "borrowed":
        raise HTTPException(status_code=409, detail="Only borrowed items can be returned.")
    inventory = db.query(InventoryItem).filter(InventoryItem.sku == item.sku).first()
    if not inventory:
        raise HTTPException(status_code=409, detail="Inventory record is missing.")
    inventory.stock_remaining += item.quantity
    item.status = "returned"
    item.reviewed_by = request.reviewer
    item.returned_at = datetime.now(timezone.utc)
    db.commit()
    db.refresh(item)
    return {"kind": "return", "message": f"{item.product} was returned and stock was updated.", "item": request_payload(item)}


@app.post("/feedback", response_model=DataCreateResponse)
def create_feedback(request: FeedbackCreate, db: Session = Depends(get_db)) -> dict:
    feedback = AgentFeedback(
        question=request.question,
        answer=request.answer,
        intent=request.intent,
        helpful=request.helpful,
        escalated=request.escalated,
        note=request.note,
        created_at=datetime.now(timezone.utc),
    )
    db.add(feedback)
    db.commit()
    db.refresh(feedback)
    message = "Support ticket created for registrar review." if feedback.escalated else "Feedback saved. Thank you."
    return {
        "kind": "ticket" if feedback.escalated else "feedback",
        "message": message,
        "item": {
            "id": feedback.id,
            "intent": feedback.intent,
            "helpful": feedback.helpful,
            "escalated": feedback.escalated,
            "status": "pending" if feedback.escalated else "reviewed",
        },
    }


@app.post("/data/inventory", response_model=DataCreateResponse)
def create_inventory_item(request: InventoryCreate, db: Session = Depends(get_db)) -> dict:
    existing = db.query(InventoryItem).filter(InventoryItem.sku == request.sku).first()
    if existing:
        raise HTTPException(status_code=409, detail="SKU already exists.")

    item = InventoryItem(
        sku=request.sku,
        product=request.product,
        category=request.category,
        stock_remaining=request.stock_remaining,
        reorder_level=request.reorder_level,
        supplier=request.supplier,
    )
    db.add(item)

    if request.note:
        db.add(
            KnowledgeDocument(
                title=f"Inventory note: {request.product}",
                source_type="AI_NOTE",
                content=request.note,
            )
        )

    db.commit()
    db.refresh(item)
    return {
        "kind": "inventory",
        "message": f"{item.product} was added to lab supply records.",
        "item": {
            "id": item.id,
            "sku": item.sku,
            "product": item.product,
            "category": item.category,
            "stock_remaining": item.stock_remaining,
            "reorder_level": item.reorder_level,
            "supplier": item.supplier,
            "note": request.note,
        },
    }


@app.post("/data/sales", response_model=DataCreateResponse)
def create_sale(request: SaleCreate, db: Session = Depends(get_db)) -> dict:
    try:
        sale_date = date.fromisoformat(request.date)
    except ValueError as exc:
        raise HTTPException(status_code=422, detail="Use date format YYYY-MM-DD.") from exc

    sale = Sale(
        date=sale_date,
        product=request.product,
        category=request.category,
        quantity=request.quantity,
        amount=request.amount,
        customer_segment=request.customer_segment,
    )
    db.add(sale)

    if request.note:
        db.add(
            KnowledgeDocument(
                title=f"Sales note: {request.product}",
                source_type="AI_NOTE",
                content=request.note,
            )
        )

    db.commit()
    db.refresh(sale)
    return {
        "kind": "sales",
        "message": f"{sale.product} was added to student activity records.",
        "item": {
            "id": sale.id,
            "date": sale.date.isoformat(),
            "product": sale.product,
            "category": sale.category,
            "quantity": sale.quantity,
            "amount": sale.amount,
            "customer_segment": sale.customer_segment,
            "note": request.note,
        },
    }


@app.post("/data/forecast-points", response_model=DataCreateResponse)
def create_forecast_point(request: ForecastPointCreate, db: Session = Depends(get_db)) -> dict:
    try:
        sale_date = date.fromisoformat(f"{request.month}-01")
    except ValueError as exc:
        raise HTTPException(status_code=422, detail="Use month format YYYY-MM.") from exc

    sale = Sale(
        date=sale_date,
        product="Forecast Planning Input",
        category="Forecast",
        quantity=1,
        amount=request.revenue,
        customer_segment="Planning",
    )
    db.add(sale)

    if request.note:
        db.add(
            KnowledgeDocument(
                title=f"Forecast note: {request.month}",
                source_type="AI_NOTE",
                content=request.note,
            )
        )

    db.commit()
    db.refresh(sale)
    return {
        "kind": "forecast",
        "message": f"Forecast point for {request.month} was added to school activity history.",
        "item": {
            "id": sale.id,
            "month": request.month,
            "revenue": sale.amount,
            "note": request.note,
        },
    }


@app.get("/documents")
def list_documents(db: Session = Depends(get_db)) -> list[dict]:
    docs = db.query(KnowledgeDocument).order_by(KnowledgeDocument.id).all()
    return [{"id": doc.id, "title": doc.title, "source_type": doc.source_type, "characters": len(doc.content)} for doc in docs]


@app.post("/documents/reindex")
def reindex_documents(db: Session = Depends(get_db)) -> dict:
    indexed_chunks = index_missing_documents(db)
    status = vector_store_status()
    return {
        "status": "completed" if status["schema_ready"] else "fallback_active",
        "indexed_chunks": indexed_chunks,
        "retrieval": status,
    }


@app.post("/documents/text", response_model=DataCreateResponse)
def create_text_document(request: DocumentCreate, db: Session = Depends(get_db)) -> dict:
    content = request.content
    if request.note:
        content = f"{content}\n\nInternal note: {request.note}"

    doc = KnowledgeDocument(title=request.title, source_type=request.source_type.upper(), content=content)
    db.add(doc)
    db.commit()
    db.refresh(doc)
    indexed_chunks = index_document(db, doc)
    return {
        "kind": "rag_document",
        "message": f"{doc.title} was added to the school document knowledge base.",
        "item": {
            "id": doc.id,
            "title": doc.title,
            "source_type": doc.source_type,
            "content": doc.content,
            "characters": len(doc.content),
            "note": request.note,
            "indexed_chunks": indexed_chunks,
        },
    }


@app.post("/ai/notes", response_model=DataCreateResponse)
def create_ai_note(request: DocumentCreate, db: Session = Depends(get_db)) -> dict:
    doc = KnowledgeDocument(
        title=request.title,
        source_type="AI_NOTE",
        content=f"{request.content}\n\nWhy it matters: {request.note or 'No extra note provided.'}",
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)
    indexed_chunks = index_document(db, doc)
    return {
        "kind": "ai_note",
        "message": f"AI note '{doc.title}' was saved and can be retrieved by RAG.",
        "item": {
            "id": doc.id,
            "title": doc.title,
            "source_type": doc.source_type,
            "content": doc.content,
            "characters": len(doc.content),
            "note": request.note,
            "indexed_chunks": indexed_chunks,
        },
    }


@app.post("/documents/upload", response_model=UploadResponse)
async def upload_document(file: UploadFile = File(...), db: Session = Depends(get_db)) -> dict:
    raw = await file.read()
    if len(raw) > 5 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File is larger than the 5 MB demo limit.")
    suffix = Path(file.filename or "").suffix.lower()
    if suffix not in {".txt", ".pdf", ".docx", ".pptx"}:
        raise HTTPException(status_code=415, detail="Supported files: TXT, PDF, DOCX, and PPTX.")

    try:
        if suffix == ".txt":
            content = raw.decode("utf-8")
        elif suffix == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(BytesIO(raw))
            content = "\n\n".join(f"Page {index}. {page.extract_text() or ''}" for index, page in enumerate(reader.pages, start=1))
        elif suffix == ".docx":
            from docx import Document

            document = Document(BytesIO(raw))
            content = "Page 1. " + "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
        else:
            from pptx import Presentation

            presentation = Presentation(BytesIO(raw))
            slides = []
            for index, slide in enumerate(presentation.slides, start=1):
                text = " ".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
                slides.append(f"Page {index}. {text}")
            content = "\n\n".join(slides)
    except Exception as exc:
        raise HTTPException(status_code=422, detail="The document could not be read or contains no extractable text.") from exc

    if not content.strip():
        raise HTTPException(status_code=422, detail="The document contains no extractable text.")

    doc = KnowledgeDocument(title=file.filename, source_type=suffix[1:].upper(), content=content)
    db.add(doc)
    db.commit()
    db.refresh(doc)
    indexed_chunks = index_document(db, doc)
    return {
        "filename": file.filename,
        "stored_characters": len(content),
        "indexed_chunks": indexed_chunks,
        "message": "Document stored and available for RAG search.",
    }
